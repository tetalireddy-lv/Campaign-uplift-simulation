"""
Build Campaign Uplift Simulator hackathon PowerPoint presentation.
Run: python3 diagrams/build_presentation.py
Output: diagrams/Campaign_Uplift_Simulator.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1E, 0x3A, 0x5F)
BLUE        = RGBColor(0x4F, 0x8E, 0xF7)
LIGHT_BLUE  = RGBColor(0xDB, 0xEA, 0xFE)
GREEN       = RGBColor(0x34, 0xC7, 0x7B)
LIGHT_GREEN = RGBColor(0xD1, 0xFA, 0xE5)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
LIGHT_AMBER = RGBColor(0xFE, 0xF3, 0xC7)
RED         = RGBColor(0xEF, 0x44, 0x44)
LIGHT_RED   = RGBColor(0xFE, 0xE2, 0xE2)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF8, 0xFA, 0xFC)
DARK_TEXT   = RGBColor(0x1E, 0x29, 0x3B)
MID_TEXT    = RGBColor(0x47, 0x55, 0x69)
LIGHT_TEXT  = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG     = RGBColor(0xF1, 0xF5, 0xF9)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

NAV_LABELS = [
    "01 Title", "02 Problem", "03 Brief Input", "04 Why",
    "05 Solution", "06 Workflow", "07 Arch", "08 Simulation",
    "09 Results", "10 Control", "11 Impact"
]
CONTENT_Y = Inches(1.18)


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=Pt(12), bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def nav_bar(slide, active_idx):
    bar_h = Inches(0.28)
    bar_y = H - bar_h
    add_rect(slide, 0, bar_y, W, bar_h, fill=NAVY)
    step_w = W / len(NAV_LABELS)
    for i, lbl in enumerate(NAV_LABELS):
        col = WHITE if i == active_idx else LIGHT_TEXT
        add_text(slide, lbl, step_w * i, bar_y, step_w, bar_h,
                 size=Pt(7.5), color=col, align=PP_ALIGN.CENTER, bold=(i == active_idx))


def header(slide, title, subtitle=None, bg=BLUE):
    hh = Inches(1.1)
    add_rect(slide, 0, 0, W, hh, fill=bg)
    add_text(slide, title, Inches(0.4), Inches(0.1), W - Inches(0.8), Inches(0.65),
             size=Pt(24), bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.72), W - Inches(0.8), Inches(0.35),
                 size=Pt(11), color=LIGHT_BLUE)


def callout_bar(slide, text, y, bg=NAVY, fg=WHITE, size=Pt(11)):
    h = Inches(0.65)
    add_rect(slide, 0, y, W, h, fill=bg)
    add_text(slide, text, Inches(0.4), y + Inches(0.1), W - Inches(0.8), Inches(0.48),
             size=size, italic=True, bold=True, color=fg, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, W, H, fill=NAVY)
add_rect(s1, 0, 0, Inches(0.08), H, fill=BLUE)

for cx, cy, cr, col in [
    (Inches(11.2), Inches(1.0), Inches(3.2), RGBColor(0x26, 0x4D, 0x7A)),
    (Inches(12.8), Inches(5.8), Inches(2.0), RGBColor(0x22, 0x43, 0x6B)),
    (Inches(10.0), Inches(6.5), Inches(1.2), RGBColor(0x20, 0x40, 0x68)),
]:
    c = s1.shapes.add_shape(9, cx - cr, cy - cr, cr * 2, cr * 2)
    c.fill.solid(); c.fill.fore_color.rgb = col
    c.line.fill.background()

add_rect(s1, Inches(0.5), Inches(1.5), Inches(2.5), Inches(0.35), fill=BLUE)
add_text(s1, "HACKATHON 2026", Inches(0.5), Inches(1.52), Inches(2.5), Inches(0.3),
         size=Pt(8.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s1, "Campaign Uplift Simulator",
         Inches(0.5), Inches(2.1), Inches(10.5), Inches(1.2),
         size=Pt(42), bold=True, color=WHITE)

add_text(s1, "AI-assisted campaign readiness, planning, and uplift simulation",
         Inches(0.5), Inches(3.4), Inches(10), Inches(0.5),
         size=Pt(16), color=LIGHT_BLUE)

add_rect(s1, Inches(0.5), Inches(4.05), Inches(3.5), Inches(0.04), fill=BLUE)

add_text(s1, '"Will this campaign move the right number — and can we prove it?"',
         Inches(0.5), Inches(4.25), Inches(10.5), Inches(0.6),
         size=Pt(13.5), italic=True, color=RGBColor(0xA5, 0xC8, 0xFF))

add_text(s1, "Team: [Your Team Name]     |     2026 Hackathon",
         Inches(0.5), Inches(6.8), Inches(7), Inches(0.4),
         size=Pt(10), color=LIGHT_TEXT)

nav_bar(s1, 0)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 2 — Problem
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, W, H, fill=OFF_WHITE)
header(s2, "Campaign Managers Are Flying Blind",
       "The brief is approved. The campaign is live. Nobody knows if it will work.")

problems = [
    ("Vague Goals",    "Briefs arrive with aspirational targets and no measurable baseline"),
    ("Weak Audience",  "Target segments are broadly defined with no behavioural signal"),
    ("No Measurement", "Success criteria are absent — no KPI, no control group, no attribution"),
    ("Rushed Handoff", "Execution begins before the plan is validated or complete"),
    ("Unproven ROI",   "Campaign ends with spend data but no proof of business impact"),
]

accent = [RED, AMBER, RED, AMBER, RED]
light  = [LIGHT_RED, LIGHT_AMBER, LIGHT_RED, LIGHT_AMBER, LIGHT_RED]

cw = Inches(2.4)
ch = Inches(3.6)
cgap = Inches(0.22)
total_cw = len(problems) * cw + (len(problems) - 1) * cgap
cx_start = (W - total_cw) / 2
cy_top = CONTENT_Y + Inches(0.18)

for i, (title, desc) in enumerate(problems):
    cx_ = cx_start + i * (cw + cgap)
    add_rect(s2, cx_, cy_top, cw, ch, fill=WHITE, line=accent[i], line_w=Pt(1.8))
    add_rect(s2, cx_, cy_top, cw, Inches(0.08), fill=accent[i])
    # Number circle
    num_c = s2.shapes.add_shape(9, cx_ + cw/2 - Inches(0.28),
                                 cy_top + Inches(0.22), Inches(0.56), Inches(0.56))
    num_c.fill.solid(); num_c.fill.fore_color.rgb = light[i]
    num_c.line.fill.background()
    add_text(s2, str(i + 1),
             cx_ + cw/2 - Inches(0.28), cy_top + Inches(0.24), Inches(0.56), Inches(0.44),
             size=Pt(14), bold=True, color=accent[i], align=PP_ALIGN.CENTER)
    add_text(s2, title,
             cx_ + Inches(0.12), cy_top + Inches(0.95), cw - Inches(0.24), Inches(0.45),
             size=Pt(12), bold=True, color=accent[i], align=PP_ALIGN.CENTER)
    add_text(s2, desc,
             cx_ + Inches(0.12), cy_top + Inches(1.5), cw - Inches(0.24), Inches(1.9),
             size=Pt(9.5), color=MID_TEXT, align=PP_ALIGN.CENTER, wrap=True)

callout_bar(s2,
    '💡  "The problem is not the campaign. The problem is the readiness of the plan behind it."',
    H - Inches(0.93), bg=LIGHT_RED, fg=RED, size=Pt(11))

nav_bar(s2, 1)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 3 — What a User Submits (Example Campaign Brief)
# ─────────────────────────────────────────────────────────────────────────────
s3b = prs.slides.add_slide(BLANK)
add_rect(s3b, 0, 0, W, H, fill=OFF_WHITE)
header(s3b, "What a Campaign Manager Actually Submits",
       "A real brief — some fields complete, many missing. The simulator finds every gap.")

# ── Left panel: brief document ─────────────────────────────────────────────
doc_x = Inches(0.3)
doc_w = Inches(7.6)
doc_y = CONTENT_Y + Inches(0.18)
doc_h = H - doc_y - Inches(0.55)
add_rect(s3b, doc_x, doc_y, doc_w, doc_h, fill=WHITE,
         line=RGBColor(0xCB, 0xD5, 0xE1), line_w=Pt(1.2))

# Brief header strip
add_rect(s3b, doc_x, doc_y, doc_w, Inches(0.46), fill=NAVY)
add_text(s3b, "📄  Campaign Brief — Q2 Product Launch",
         doc_x + Inches(0.2), doc_y + Inches(0.08), doc_w - Inches(0.3), Inches(0.34),
         size=Pt(12), bold=True, color=WHITE)

# Brief fields: (label, value, status)  status: "ok" | "warn" | "missing"
brief_fields = [
    ("Campaign Name",      "Q2 Premium Product Launch — New SKU",       "ok"),
    ("Objective",          "Drive awareness and trial among premium segment", "ok"),
    ("Target Audience",    "Adults 25–45, household income >$75k",       "ok"),
    ("Channels",           "Digital, Social Media, Out-of-Home",         "ok"),
    ("Budget",             "$250,000 total",                             "ok"),
    ("Timeline",           "6 weeks starting June 2026",                 "ok"),
    ("Primary KPI",        "Increase sales",                             "warn"),
    ("KPI Baseline",       "Not provided",                               "missing"),
    ("Success Threshold",  "Not defined",                                "missing"),
    ("Control Strategy",   "Not defined",                                "missing"),
    ("Attribution Method", "Mentioned but unspecified",                  "warn"),
    ("Audience Segments",  "Broad — no behavioural split provided",      "warn"),
    ("ROI Expectation",    "Not provided",                               "missing"),
]

status_cfg = {
    "ok":      (GREEN,  LIGHT_GREEN,  "✓"),
    "warn":    (AMBER,  LIGHT_AMBER,  "⚠"),
    "missing": (RED,    LIGHT_RED,    "✗"),
}

row_h  = Inches(0.355)
row_gap = Inches(0.025)
field_y = doc_y + Inches(0.55)

for i, (lbl, val, status) in enumerate(brief_fields):
    col, lcol, icon = status_cfg[status]
    ry = field_y + i * (row_h + row_gap)
    bg = lcol if status != "ok" else WHITE
    add_rect(s3b, doc_x, ry, doc_w, row_h, fill=bg)
    # thin left border
    add_rect(s3b, doc_x, ry, Inches(0.055), row_h, fill=col)
    # status icon box
    add_rect(s3b, doc_x + Inches(0.1), ry + Inches(0.05),
             Inches(0.3), row_h - Inches(0.1), fill=col)
    add_text(s3b, icon,
             doc_x + Inches(0.1), ry + Inches(0.06),
             Inches(0.3), row_h - Inches(0.1),
             size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # label
    add_text(s3b, lbl,
             doc_x + Inches(0.5), ry + Inches(0.07), Inches(1.9), row_h - Inches(0.08),
             size=Pt(9), bold=True, color=DARK_TEXT)
    # value
    add_text(s3b, val,
             doc_x + Inches(2.5), ry + Inches(0.07), Inches(4.95), row_h - Inches(0.08),
             size=Pt(9), color=(RED if status == "missing" else
                                AMBER if status == "warn" else MID_TEXT),
             italic=(status != "ok"), bold=(status == "missing"))
    # subtle row divider
    add_rect(s3b, doc_x, ry + row_h, doc_w, Inches(0.01),
             fill=RGBColor(0xE2, 0xE8, 0xF0))

# ── Right panel: readiness score + gap summary ─────────────────────────────
rp_x = doc_x + doc_w + Inches(0.3)
rp_w = W - rp_x - Inches(0.25)
rp_y = doc_y

# Readiness score card
add_rect(s3b, rp_x, rp_y, rp_w, Inches(1.6), fill=NAVY)
add_text(s3b, "Readiness Score",
         rp_x + Inches(0.2), rp_y + Inches(0.15), rp_w - Inches(0.3), Inches(0.3),
         size=Pt(10), bold=True, color=LIGHT_BLUE)
add_text(s3b, "38 / 100",
         rp_x + Inches(0.2), rp_y + Inches(0.45), rp_w - Inches(0.3), Inches(0.75),
         size=Pt(34), bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s3b, "NOT READY FOR LAUNCH",
         rp_x + Inches(0.2), rp_y + Inches(1.22), rp_w - Inches(0.3), Inches(0.28),
         size=Pt(8.5), bold=True, color=LIGHT_RED, align=PP_ALIGN.CENTER)

# Gap summary cards
gap_items = [
    ("4", "Critical\nGaps",   RED,   LIGHT_RED),
    ("3", "Warnings",         AMBER, LIGHT_AMBER),
    ("6", "Fields\nComplete", GREEN, LIGHT_GREEN),
]
gi_h = Inches(1.0)
gi_gap = Inches(0.14)
gi_y = rp_y + Inches(1.72)
gi_w = (rp_w - 2 * gi_gap) / 3

for j, (num, label, col, lcol) in enumerate(gap_items):
    gx_ = rp_x + j * (gi_w + gi_gap)
    add_rect(s3b, gx_, gi_y, gi_w, gi_h, fill=lcol, line=col, line_w=Pt(1.5))
    add_text(s3b, num,
             gx_, gi_y + Inches(0.08), gi_w, Inches(0.55),
             size=Pt(26), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s3b, label,
             gx_, gi_y + Inches(0.6), gi_w, Inches(0.38),
             size=Pt(8), color=col, align=PP_ALIGN.CENTER, bold=True)

# What the tool does next
next_y = gi_y + gi_h + Inches(0.2)
add_rect(s3b, rp_x, next_y, rp_w, Inches(0.32), fill=BLUE)
add_text(s3b, "  What happens next:",
         rp_x, next_y + Inches(0.05), rp_w, Inches(0.25),
         size=Pt(9), bold=True, color=WHITE)

next_steps = [
    "Asks 4 targeted clarification questions",
    "Proposes assumptions for missing fields",
    "Awaits manager approval before proceeding",
    "Builds plan only after sign-off",
]
for k, ns in enumerate(next_steps):
    nsy = next_y + Inches(0.38) + k * Inches(0.5)
    add_rect(s3b, rp_x, nsy, rp_w, Inches(0.44),
             fill=LIGHT_BLUE if k % 2 == 0 else WHITE,
             line=BLUE, line_w=Pt(0.5))
    add_text(s3b, f"→  {ns}",
             rp_x + Inches(0.12), nsy + Inches(0.08), rp_w - Inches(0.15), Inches(0.3),
             size=Pt(9), color=DARK_TEXT)

# Legend
leg_y3b = H - Inches(0.82)
for li, (icon, lbl, col) in enumerate([("✓", "Complete", GREEN),
                                        ("⚠", "Warning / Vague", AMBER),
                                        ("✗", "Missing / Critical", RED)]):
    lx_ = doc_x + li * Inches(2.4)
    add_rect(s3b, lx_, leg_y3b, Inches(0.26), Inches(0.26), fill=col)
    add_text(s3b, icon, lx_, leg_y3b + Inches(0.03), Inches(0.26), Inches(0.22),
             size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s3b, lbl, lx_ + Inches(0.32), leg_y3b + Inches(0.04), Inches(2.0), Inches(0.22),
             size=Pt(9), color=MID_TEXT)

nav_bar(s3b, 2)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 4 — Why This Matters
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, W, H, fill=OFF_WHITE)
header(s3, "The Cost of an Unready Campaign Brief",
       "Every gap in the brief becomes a gap in your results.")

impacts = [
    ("💸  Wasted Budget",
     "Teams commit spend without expected uplift visibility. Budget decisions are made on instinct, not data.",
     RED, LIGHT_RED),
    ("📉  Lost Attribution",
     "ROI is guessed after the fact. There is no way to separate the campaign's effect from market noise.",
     AMBER, LIGHT_AMBER),
    ("⏳  Slow Planning",
     "Teams spend 30–50% of planning time clarifying basic inputs that should have been in the original brief.",
     AMBER, LIGHT_AMBER),
    ("🕳  Manager Blindspot",
     "Campaign managers lose control to black-box tools that don't explain their recommendations or outputs.",
     RED, LIGHT_RED),
]

iw   = Inches(6.1)
ih   = Inches(2.35)
igap = Inches(0.22)
xs_  = [Inches(0.3), Inches(6.73)]
ys_  = [CONTENT_Y + Inches(0.15), CONTENT_Y + ih + igap + Inches(0.15)]

for idx, (title, desc, col, lcol) in enumerate(impacts):
    xi_ = xs_[idx % 2]
    yi_ = ys_[idx // 2]
    add_rect(s3, xi_, yi_, iw, ih, fill=WHITE, line=col, line_w=Pt(2))
    add_rect(s3, xi_, yi_, Inches(0.09), ih, fill=col)
    add_text(s3, title, xi_ + Inches(0.22), yi_ + Inches(0.2), iw - Inches(0.3), Inches(0.42),
             size=Pt(14), bold=True, color=col)
    add_text(s3, desc, xi_ + Inches(0.22), yi_ + Inches(0.72), iw - Inches(0.3), Inches(1.5),
             size=Pt(10.5), color=MID_TEXT, wrap=True)

callout_bar(s3, '"Campaigns should be planned around proof — not hope."',
            H - Inches(0.93))

nav_bar(s3, 3)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 4 — Solution Overview
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, W, H, fill=OFF_WHITE)
header(s4, "Campaign Uplift Simulator — What It Does",
       "Eight structured steps from vague brief to validated campaign plan with uplift forecast.")

steps = [
    ("01", "Parse Brief",         NAVY),
    ("02", "Score Readiness",     BLUE),
    ("03", "Flag Gaps",           RED),
    ("04", "Clarify",             AMBER),
    ("05", "Approve Assumptions", GREEN),
    ("06", "Build Plan",          BLUE),
    ("07", "Simulate Uplift",     NAVY),
    ("08", "Final Handoff",       GREEN),
]

step_w  = Inches(1.44)
step_h  = Inches(0.78)
arrow_w = Inches(0.22)
total_sw = len(steps) * step_w + (len(steps) - 1) * arrow_w
sx_start = (W - total_sw) / 2
sy_      = CONTENT_Y + Inches(0.3)

for i, (num, label, col) in enumerate(steps):
    sxi = sx_start + i * (step_w + arrow_w)
    add_rect(s4, sxi, sy_, step_w, step_h, fill=col)
    add_text(s4, num, sxi, sy_ + Inches(0.07), step_w, Inches(0.24),
             size=Pt(8), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, label, sxi, sy_ + Inches(0.35), step_w, Inches(0.38),
             size=Pt(8.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(s4, "▶", sxi + step_w, sy_ + Inches(0.24), arrow_w, Inches(0.3),
                 size=Pt(9), color=NAVY, align=PP_ALIGN.CENTER)

# Three value pillars
pillars = [
    ("🏗  Structured Readiness",
     "Turns vague, incomplete briefs into validated, schema-driven campaign plans ready for execution.",
     BLUE, LIGHT_BLUE),
    ("👤  Human-in-the-Loop",
     "Manager approves every AI-generated assumption before the workflow advances to the next step.",
     GREEN, LIGHT_GREEN),
    ("📊  Deterministic Simulation",
     "Same inputs always produce the same uplift estimate. Fully auditable, traceable, and explainable.",
     NAVY, LIGHT_BLUE),
]

pw   = Inches(4.1)
ph   = Inches(2.5)
pgap = Inches(0.22)
total_pw = len(pillars) * pw + (len(pillars) - 1) * pgap
px_start_ = (W - total_pw) / 2
py_ = sy_ + step_h + Inches(0.45)

for i, (ptitle, pdesc, pcol, plcol) in enumerate(pillars):
    pxi = px_start_ + i * (pw + pgap)
    add_rect(s4, pxi, py_, pw, ph, fill=plcol, line=pcol, line_w=Pt(1.8))
    add_rect(s4, pxi, py_, pw, Inches(0.07), fill=pcol)
    add_text(s4, ptitle, pxi + Inches(0.2), py_ + Inches(0.2), pw - Inches(0.3), Inches(0.45),
             size=Pt(13), bold=True, color=pcol)
    add_text(s4, pdesc, pxi + Inches(0.2), py_ + Inches(0.78), pw - Inches(0.3), Inches(1.6),
             size=Pt(10.5), color=DARK_TEXT, wrap=True)

callout_bar(s4,
    '"AI structures and validates. The manager decides. The simulation predicts. Everyone sees the same numbers before launch."',
    H - Inches(0.93))

nav_bar(s4, 4)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 5 — Workflow / User Journey
# ─────────────────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, W, H, fill=OFF_WHITE)
header(s5, "Human-in-the-Loop, Every Step of the Way",
       "AI structures and suggests. The manager validates, approves, and controls.")

flow_steps = [
    ("Brief Input",              BLUE,  False),
    ("Understand Brief",         BLUE,  False),
    ("Readiness Score",          BLUE,  False),
    ("Gap Detection",            RED,   False),
    ("Clarification Questions",  AMBER, True),
    ("Assumption Approval",      GREEN, True),
    ("Campaign Plan Generated",  BLUE,  False),
    ("Uplift Simulation",        NAVY,  False),
    ("QA Check",                 BLUE,  False),
    ("Final Handoff",            GREEN, True),
]

n      = len(flow_steps)
node_w = Inches(1.15)
node_h = Inches(0.58)
h_gap  = Inches(0.09)
total_fw = n * node_w + (n - 1) * h_gap
fx_s = (W - total_fw) / 2
fy_  = CONTENT_Y + Inches(0.35)

for i, (label, col, is_human) in enumerate(flow_steps):
    fx_ = fx_s + i * (node_w + h_gap)
    nh  = node_h + (Inches(0.18) if is_human else 0)
    ny_ = fy_   - (Inches(0.09)  if is_human else 0)
    bw_ = Pt(2) if is_human else Pt(0)
    add_rect(s5, fx_, ny_, node_w, nh, fill=col,
             line=(WHITE if is_human else None), line_w=bw_)
    add_text(s5, label, fx_, ny_ + Inches(0.1), node_w, nh - Inches(0.1),
             size=Pt(7.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if is_human:
        add_text(s5, "👤 Manager",
                 fx_, ny_ + nh, node_w, Inches(0.22),
                 size=Pt(6.5), bold=True, color=col, align=PP_ALIGN.CENTER)
    if i < n - 1:
        add_text(s5, "▶", fx_ + node_w, fy_ + Inches(0.14), h_gap, Inches(0.3),
                 size=Pt(7), color=NAVY, align=PP_ALIGN.CENTER)

# Legend
leg_y = fy_ + Inches(1.02)
add_rect(s5, Inches(0.3), leg_y, Inches(0.22), Inches(0.22), fill=BLUE)
add_text(s5, "AI Step", Inches(0.57), leg_y, Inches(1.5), Inches(0.22),
         size=Pt(9), color=MID_TEXT)
add_rect(s5, Inches(2.2), leg_y, Inches(0.22), Inches(0.22), fill=GREEN)
add_text(s5, "Manager Approval Required", Inches(2.47), leg_y, Inches(3.0), Inches(0.22),
         size=Pt(9), color=MID_TEXT)

# Control points — full width 3-column
control_pts = [
    ("✓  Approve or reject any assumption",    GREEN),
    ("✓  Edit channel mix before simulation",  GREEN),
    ("✓  Sign off before final handoff",        GREEN),
    ("✓  Nothing advances without approval",   GREEN),
    ("✓  All assumptions visible and editable",GREEN),
    ("✓  Every uplift estimate is explainable",GREEN),
]
cp_w   = Inches(4.0)
cp_h   = Inches(0.52)
cp_gap = Inches(0.16)
cp_y   = leg_y + Inches(0.42)
cp_xs  = [Inches(0.3), Inches(4.62), Inches(8.94)]

for ci, (cpt, ccol) in enumerate(control_pts):
    cxi_ = cp_xs[ci % 3]
    cyi_ = cp_y + (ci // 3) * (cp_h + cp_gap)
    add_rect(s5, cxi_, cyi_, cp_w, cp_h, fill=LIGHT_GREEN, line=GREEN, line_w=Pt(1))
    add_text(s5, cpt, cxi_ + Inches(0.12), cyi_ + Inches(0.1), cp_w - Inches(0.15), Inches(0.35),
             size=Pt(9.5), color=DARK_TEXT)

callout_bar(s5,
    '"This is not automation. This is assisted decision-making with the manager in control."',
    H - Inches(0.93))

nav_bar(s5, 5)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 6 — Architecture
# ─────────────────────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, W, H, fill=RGBColor(0x0F, 0x17, 0x2A))
header(s6, "Built for Reliability, Designed for Extensibility",
       "Production-grade architecture. Hackathon-ready execution.", bg=NAVY)

arrow_col = RGBColor(0x4B, 0x5E, 0x82)

layers = [
    ("React Frontend  (TypeScript)",       BLUE),
    ("FastAPI Backend  (Python)",          RGBColor(0x7C, 0x3A, 0xED)),
    ("LangGraph Orchestrated Workflow",    RGBColor(0xEC, 0x48, 0x99)),
]

lw   = Inches(4.5)
lh   = Inches(0.52)
lgap = Inches(0.2)
lx_  = Inches(0.4)
ly_  = CONTENT_Y + Inches(0.18)

for i, (lbl, col) in enumerate(layers):
    add_rect(s6, lx_, ly_ + i * (lh + lgap), lw, lh, fill=col)
    add_text(s6, lbl, lx_, ly_ + i * (lh + lgap) + Inches(0.1), lw, lh - Inches(0.1),
             size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(layers) - 1:
        ay_ = ly_ + (i + 1) * (lh + lgap) - lgap
        add_text(s6, "▼", lx_ + lw / 2 - Inches(0.15), ay_, Inches(0.3), lgap,
                 size=Pt(9), color=arrow_col, align=PP_ALIGN.CENTER)

# Agent grid (2 rows × 4 cols)
agents = [
    ("Brief Parser",        BLUE),
    ("Readiness Validator", BLUE),
    ("Gap Detector",        RED),
    ("Clarification Agent", AMBER),
    ("Planning Agent",      GREEN),
    ("Simulation Layer",    RGBColor(0xEC, 0x48, 0x99)),
    ("QA Agent",            NAVY),
    ("Handoff Generator",   GREEN),
]
aw_   = Inches(1.0)
ah_   = Inches(0.48)
agap_ = Inches(0.1)
agent_y = ly_ + len(layers) * (lh + lgap) + Inches(0.14)
ax_s_ = lx_ + Inches(0.1)

for i, (albl, acol) in enumerate(agents):
    row = i // 4; col_i = i % 4
    axi_ = ax_s_ + col_i * (aw_ + agap_)
    ayi_ = agent_y + row * (ah_ + Inches(0.1))
    add_rect(s6, axi_, ayi_, aw_, ah_, fill=acol)
    add_text(s6, albl, axi_, ayi_ + Inches(0.08), aw_, ah_ - Inches(0.08),
             size=Pt(8), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

arr2_y = agent_y + 2 * (ah_ + Inches(0.1))
add_text(s6, "▼", lx_ + lw / 2 - Inches(0.15), arr2_y, Inches(0.3), Inches(0.26),
         size=Pt(10), color=arrow_col, align=PP_ALIGN.CENTER)

data_y = arr2_y + Inches(0.3)
for j, (dlbl, dcol) in enumerate([
    ("CSV Benchmark Data",        RGBColor(0x0E, 0x74, 0x90)),
    ("Manager Assumptions",       RGBColor(0x06, 0x78, 0x50)),
    ("Final Plan + Uplift Output",RGBColor(0x7C, 0x3A, 0xED)),
]):
    ddw = Inches(1.38)
    ddx = lx_ + Inches(0.05) + j * (ddw + Inches(0.07))
    add_rect(s6, ddx, data_y, ddw, Inches(0.48), fill=dcol)
    add_text(s6, dlbl, ddx, data_y + Inches(0.1), ddw, Inches(0.35),
             size=Pt(8), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Stack highlights (right panel)
hx_ = Inches(5.35)
hy_ = CONTENT_Y + Inches(0.18)
add_rect(s6, hx_, hy_, Inches(7.6), Inches(5.7),
         fill=RGBColor(0x16, 0x24, 0x3B), line=BLUE, line_w=Pt(1.2))
add_text(s6, "Stack Highlights",
         hx_ + Inches(0.25), hy_ + Inches(0.2), Inches(7.0), Inches(0.35),
         size=Pt(12), bold=True, color=BLUE)

highlights = [
    ("⚡", "Async / parallel LLM calls where beneficial"),
    ("📊", "CSV-backed benchmarks for deterministic simulation"),
    ("🔄", "LangGraph multi-agent orchestration"),
    ("🛡", "Mock fallback layer — works without live credentials"),
    ("🔗", "FastAPI async endpoints"),
    ("⚛️",  "React + TypeScript frontend"),
    ("🧪", "Deterministic simulation — same inputs, same output, always"),
]
for k, (icon, htxt) in enumerate(highlights):
    row = k // 2; col_ = k % 2
    hxi_ = hx_ + Inches(0.25) + col_ * Inches(3.65)
    hyi_ = hy_ + Inches(0.75) + row * Inches(0.88)
    add_rect(s6, hxi_, hyi_, Inches(3.5), Inches(0.72),
             fill=RGBColor(0x1E, 0x35, 0x58), line=BLUE, line_w=Pt(0.5))
    add_text(s6, icon, hxi_ + Inches(0.1), hyi_ + Inches(0.12),
             Inches(0.38), Inches(0.45), size=Pt(14), align=PP_ALIGN.CENTER, color=WHITE)
    add_text(s6, htxt, hxi_ + Inches(0.55), hyi_ + Inches(0.14),
             Inches(2.85), Inches(0.48), size=Pt(9), color=WHITE, wrap=True)

nav_bar(s6, 6)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 7 — Market Uplift Simulation
# ─────────────────────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, W, H, fill=OFF_WHITE)
header(s7, "Three Scenarios. One Decision.",
       "Uplift is not a guess — it is a function of what you put in.")

scenarios = [
    ("Scenario 1", "Current Brief As-Is",
     "Uses only what's in the brief today. All gaps and missing inputs remain as-is.",
     "Low Confidence", RED, LIGHT_RED),
    ("Scenario 2", "Clarified Plan",
     "Applies manager-approved assumptions to fill the identified gaps in the brief.",
     "Medium Confidence", AMBER, LIGHT_AMBER),
    ("Scenario 3", "Optimized Plan",
     "Adds improved targeting, better channel mix, and higher execution quality.",
     "High Confidence", GREEN, LIGHT_GREEN),
]

sw_   = Inches(4.05)
sh_   = Inches(2.1)
sgap_ = Inches(0.26)
total_sw2 = len(scenarios) * sw_ + (len(scenarios) - 1) * sgap_
sx_s2 = (W - total_sw2) / 2
sy_2  = CONTENT_Y + Inches(0.2)

for i, (badge, title, desc, conf, col, lcol) in enumerate(scenarios):
    sxi_ = sx_s2 + i * (sw_ + sgap_)
    add_rect(s7, sxi_, sy_2, sw_, sh_, fill=WHITE, line=col, line_w=Pt(2))
    add_rect(s7, sxi_, sy_2, sw_, Inches(0.07), fill=col)
    add_rect(s7, sxi_ + Inches(0.15), sy_2 + Inches(0.16), Inches(1.15), Inches(0.28), fill=lcol)
    add_text(s7, badge, sxi_ + Inches(0.15), sy_2 + Inches(0.16), Inches(1.15), Inches(0.28),
             size=Pt(7.5), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s7, title, sxi_ + Inches(0.15), sy_2 + Inches(0.54), sw_ - Inches(0.25), Inches(0.42),
             size=Pt(13), bold=True, color=col)
    add_text(s7, desc, sxi_ + Inches(0.15), sy_2 + Inches(1.05), sw_ - Inches(0.25), Inches(0.78),
             size=Pt(9.5), color=MID_TEXT, wrap=True)
    add_text(s7, f"● {conf}",
             sxi_ + Inches(0.15), sy_2 + sh_ - Inches(0.35), sw_ - Inches(0.25), Inches(0.28),
             size=Pt(9), bold=True, color=col)

# Formula box
fy2_ = sy_2 + sh_ + Inches(0.28)
add_rect(s7, Inches(0.3), fy2_, W - Inches(0.6), Inches(0.82), fill=NAVY, line=BLUE, line_w=Pt(1))
add_text(s7,
    "Uplift  =  f ( Historical Benchmarks,  Audience Segment Mix,  Channel Fit Score,"
    "  Budget Strength,  Message Clarity,  Execution Confidence )",
    Inches(0.5), fy2_ + Inches(0.12), W - Inches(1.0), Inches(0.6),
    size=Pt(10.5), italic=True, color=WHITE, align=PP_ALIGN.CENTER)

# Audience segments — 4 columns full width
segments = [
    ("Persuadables",   "High-value — respond to messaging. Focus spend here.", GREEN),
    ("Sure Things",    "Will convert regardless — don't overspend.",           BLUE),
    ("Lost Causes",    "Won't convert — exclude from targeting.",              RED),
    ("Do Not Disturb", "Risk opt-out — handle carefully.",                     AMBER),
]
seg_w_ = Inches(3.02)
seg_h_ = Inches(0.88)
seg_gap_ = Inches(0.22)
total_segw = len(segments) * seg_w_ + (len(segments) - 1) * seg_gap_
seg_x_s = (W - total_segw) / 2
seg_y_ = fy2_ + Inches(1.0)

for i, (sname, sdesc, scol) in enumerate(segments):
    sxi_ = seg_x_s + i * (seg_w_ + seg_gap_)
    add_rect(s7, sxi_, seg_y_, seg_w_, seg_h_, fill=WHITE, line=scol, line_w=Pt(1.5))
    add_rect(s7, sxi_, seg_y_, Inches(0.07), seg_h_, fill=scol)
    add_text(s7, sname, sxi_ + Inches(0.18), seg_y_ + Inches(0.08),
             seg_w_ - Inches(0.22), Inches(0.32), size=Pt(10.5), bold=True, color=scol)
    add_text(s7, sdesc, sxi_ + Inches(0.18), seg_y_ + Inches(0.44),
             seg_w_ - Inches(0.22), Inches(0.38), size=Pt(8.5), color=MID_TEXT, wrap=True)

nav_bar(s7, 7)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 8 — Results / Example Output
# ─────────────────────────────────────────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, W, H, fill=OFF_WHITE)
header(s8, "Before and After — What the Simulator Shows You",
       "Concrete numbers before launch. Not guesses after.")

# Results table
headers8 = ["Scenario", "Expected KPI", "Market Uplift", "ROI Estimate", "Confidence"]
rows8 = [
    ["Current Brief As-Is", "1.2% conversion", "+4%",  "1.1×", "Low"],
    ["Clarified Plan",       "2.4% conversion", "+11%", "2.3×", "Medium"],
    ["Optimized Plan",       "3.8% conversion", "+19%", "3.6×", "High"],
]
row_cols8 = [LIGHT_RED, LIGHT_AMBER, LIGHT_GREEN]
conf_cs   = [RED, AMBER, GREEN]
col_w8    = [Inches(2.8), Inches(2.1), Inches(1.9), Inches(1.9), Inches(1.6)]
tx8 = Inches(0.3)
ty8 = CONTENT_Y + Inches(0.15)
th8 = Inches(0.45)
tw8 = sum(col_w8)

add_rect(s8, tx8, ty8, tw8, th8, fill=NAVY)
cxr = tx8
for i, hdr in enumerate(headers8):
    add_text(s8, hdr, cxr + Inches(0.04), ty8 + Inches(0.08),
             col_w8[i] - Inches(0.08), th8 - Inches(0.1),
             size=Pt(10.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    cxr += col_w8[i]

for r, (rdata, rcol, ccol) in enumerate(zip(rows8, row_cols8, conf_cs)):
    ry8 = ty8 + (r + 1) * (th8 + Inches(0.04))
    add_rect(s8, tx8, ry8, tw8, th8, fill=rcol, line=WHITE, line_w=Pt(0.8))
    cxr = tx8
    for c, cell in enumerate(rdata):
        fc_ = ccol if c == 4 else DARK_TEXT
        add_text(s8, cell, cxr + Inches(0.04), ry8 + Inches(0.08),
                 col_w8[c] - Inches(0.08), th8 - Inches(0.1),
                 size=Pt(10.5), bold=(c == 4), color=fc_, align=PP_ALIGN.CENTER)
        cxr += col_w8[c]

# Bar chart beside table
chart_x8 = tx8 + tw8 + Inches(0.35)
chart_y8 = ty8
ch8_     = Inches(1.9)
bar_vs   = [4, 11, 19]
bar_cs_  = [RED, AMBER, GREEN]
bar_ls   = ["As-Is", "Clarified", "Optimized"]
max_v    = 20
bw8      = Inches(0.75)
b_gap8   = Inches(0.25)
bx_s8    = chart_x8 + Inches(0.2)

add_text(s8, "Market Uplift %", chart_x8, chart_y8, Inches(3.5), Inches(0.3),
         size=Pt(9), bold=True, color=MID_TEXT, align=PP_ALIGN.CENTER)
for bi, (bv, bc, bl) in enumerate(zip(bar_vs, bar_cs_, bar_ls)):
    bxi_ = bx_s8 + bi * (bw8 + b_gap8)
    bh__ = ch8_ * (bv / max_v)
    byi_ = chart_y8 + ch8_ - bh__ + Inches(0.32)
    add_rect(s8, bxi_, byi_, bw8, bh__, fill=bc)
    add_text(s8, f"+{bv}%", bxi_, byi_ - Inches(0.26), bw8, Inches(0.24),
             size=Pt(9.5), bold=True, color=bc, align=PP_ALIGN.CENTER)
    add_text(s8, bl, bxi_, chart_y8 + ch8_ + Inches(0.35), bw8, Inches(0.28),
             size=Pt(8.5), color=MID_TEXT, align=PP_ALIGN.CENTER)

# Insight cards — full width, 2 columns
insight_y8 = ty8 + 4 * (th8 + Inches(0.04)) + Inches(0.22)
insights8 = [
    ("Missing baseline reduces simulation confidence by ~40%",  RED),
    ("Clarified audience segment improves expected conversion",  GREEN),
    ("Optimized channel mix improves ROI 56% over as-is plan",  GREEN),
    ("Control strategy improves attribution measurability",      BLUE),
    ("Clear budget allocation increases simulation confidence",  BLUE),
    ("High message clarity score adds ~8% to uplift estimate",   AMBER),
]
iw8   = Inches(6.25)
ih8   = Inches(0.52)
ig8   = Inches(0.13)
ix8_s = [Inches(0.3), Inches(6.73)]

for ii, (itxt, icol) in enumerate(insights8):
    col_ = ii % 2; row_ = ii // 2
    ixi_ = ix8_s[col_]
    iyi_ = insight_y8 + row_ * (ih8 + ig8)
    add_rect(s8, ixi_, iyi_, iw8, ih8, fill=WHITE, line=icol, line_w=Pt(1.5))
    add_rect(s8, ixi_, iyi_, Inches(0.07), ih8, fill=icol)
    add_text(s8, f"💡  {itxt}",
             ixi_ + Inches(0.2), iyi_ + Inches(0.1), iw8 - Inches(0.25), ih8 - Inches(0.1),
             size=Pt(9.5), color=DARK_TEXT)

nav_bar(s8, 8)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 9 — Manager Control & Safety
# ─────────────────────────────────────────────────────────────────────────────
s9 = prs.slides.add_slide(BLANK)
add_rect(s9, 0, 0, W, H, fill=OFF_WHITE)
header(s9, "AI Assists. The Manager Decides.",
       "Transparent, auditable, and fully manager-controlled — by design.")

controls9 = [
    ("🚫", "No Auto-Launch",        "AI never triggers a campaign action. Launch is always a manual step."),
    ("✅", "Assumption Approval",   "Manager reviews and signs off on every AI-generated assumption."),
    ("✏️", "Editable Channel Mix",   "Manager can adjust the channel mix before the simulation runs."),
    ("🔍", "Explainable Estimates",  "Every uplift number has a visible, traceable reason attached to it."),
    ("📋", "Visible Assumptions",    "All assumptions are shown in the assumption register — always visible."),
    ("🔒", "Manager Sign-Off",       "Final handoff packet is released only after explicit manager approval."),
]

cw9  = Inches(4.0)
ch9  = Inches(1.65)
cg9  = Inches(0.22)
cy9_ = CONTENT_Y + Inches(0.22)
cx9s = [Inches(0.3), Inches(4.62), Inches(8.94)]

for i, (icon, ctitle, cdesc) in enumerate(controls9):
    col_i = i % 3; row_i = i // 3
    cxi_ = cx9s[col_i]
    cyi_ = cy9_ + row_i * (ch9 + cg9)
    add_rect(s9, cxi_, cyi_, cw9, ch9, fill=WHITE, line=GREEN, line_w=Pt(1.8))
    add_rect(s9, cxi_, cyi_, cw9, Inches(0.07), fill=GREEN)
    ic9c = s9.shapes.add_shape(9, cxi_ + Inches(0.16), cyi_ + Inches(0.2),
                               Inches(0.5), Inches(0.5))
    ic9c.fill.solid(); ic9c.fill.fore_color.rgb = LIGHT_GREEN
    ic9c.line.fill.background()
    add_text(s9, icon, cxi_ + Inches(0.16), cyi_ + Inches(0.2),
             Inches(0.5), Inches(0.45), size=Pt(14), align=PP_ALIGN.CENTER, color=GREEN)
    add_text(s9, ctitle, cxi_ + Inches(0.8), cyi_ + Inches(0.24),
             cw9 - Inches(0.9), Inches(0.4), size=Pt(12), bold=True, color=GREEN)
    add_text(s9, cdesc, cxi_ + Inches(0.18), cyi_ + Inches(0.78),
             cw9 - Inches(0.28), Inches(0.78), size=Pt(9.5), color=MID_TEXT, wrap=True)

callout_bar(s9,
    '"This tool is not here to replace judgment — it is here to make the manager\'s judgment better informed."',
    H - Inches(0.93), bg=RGBColor(0x05, 0x2E, 0x16))

nav_bar(s9, 9)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 10 — Impact & Future Scope
# ─────────────────────────────────────────────────────────────────────────────
s10 = prs.slides.add_slide(BLANK)
add_rect(s10, 0, 0, W, H, fill=OFF_WHITE)
header(s10, "Faster Plans. Better Campaigns. Measurable Outcomes.",
       "What Campaign Uplift Simulator delivers today — and where it goes next.")

# Impact column header
imp_x = Inches(0.3); imp_w = Inches(5.9)
add_rect(s10, imp_x, CONTENT_Y + Inches(0.12), imp_w, Inches(0.4), fill=GREEN)
add_text(s10, "  Impact — Today",
         imp_x, CONTENT_Y + Inches(0.12), imp_w, Inches(0.4),
         size=Pt(12), bold=True, color=WHITE)

impact_items10 = [
    "Brief-to-plan conversion time reduced significantly",
    "Campaign readiness visible before launch, not after",
    "ROI expectations set before budget is committed",
    "Handoff quality improved across campaign teams",
    "Manager confidence increased at every planning step",
]
for ii, item in enumerate(impact_items10):
    iy_ = CONTENT_Y + Inches(0.68) + ii * Inches(0.7)
    add_rect(s10, imp_x, iy_, imp_w, Inches(0.58), fill=LIGHT_GREEN, line=GREEN, line_w=Pt(1))
    add_text(s10, f"✓  {item}",
             imp_x + Inches(0.15), iy_ + Inches(0.11), imp_w - Inches(0.2), Inches(0.38),
             size=Pt(10), color=DARK_TEXT)

# Future scope column header
fut_x = Inches(6.8); fut_w = Inches(6.2)
add_rect(s10, fut_x, CONTENT_Y + Inches(0.12), fut_w, Inches(0.4), fill=BLUE)
add_text(s10, "  Future Scope",
         fut_x, CONTENT_Y + Inches(0.12), fut_w, Inches(0.4),
         size=Pt(12), bold=True, color=WHITE)

future_items10 = [
    "CRM integration for real audience data",
    "Live benchmark feeds from ad platforms",
    "A/B testing recommendation engine",
    "Budget allocation optimizer",
    "Multi-touch attribution modeling",
    "Auto-generated creative briefs",
]
for fi, item in enumerate(future_items10):
    fy_ = CONTENT_Y + Inches(0.68) + fi * Inches(0.58)
    add_rect(s10, fut_x, fy_, fut_w, Inches(0.48), fill=LIGHT_BLUE, line=BLUE, line_w=Pt(1))
    add_text(s10, f"→  {item}",
             fut_x + Inches(0.15), fy_ + Inches(0.09), fut_w - Inches(0.2), Inches(0.32),
             size=Pt(10), color=DARK_TEXT)

# Closing callout — full bleed, double height for impact
callout_bar(s10,
    '"Campaign Uplift Simulator shifts planning from \'Did we launch?\''
    ' to \'Will this campaign move the right number — and can we prove it?\'"',
    H - Inches(0.93))

nav_bar(s10, 10)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/home/v-tetreddy/Campaign-uplift-simulation/diagrams/Campaign_Uplift_Simulator.pptx"
prs.save(out)
print(f"Saved: {out}")
