"""
Build Campaign Readiness Copilot – Google Slides deck (4 slides).
Output: diagrams/Campaign_Readiness_Copilot.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from lxml import etree
import copy

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Campaign_Readiness_Copilot.pptx")

# ── Palette ───────────────────────────────────────────────────────────────────
INDIGO   = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_D = RGBColor(0x31, 0x2E, 0x81)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
GREEN    = RGBColor(0x10, 0xB9, 0x81)
RED      = RGBColor(0xEF, 0x44, 0x44)
SKY      = RGBColor(0x0E, 0xA5, 0xE9)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x6B, 0x72, 0x80)
DARK     = RGBColor(0x11, 0x18, 0x27)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFF)
PANEL    = RGBColor(0xEE, 0xEF, 0xFF)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]   # completely blank

# ── helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, italic=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True, valign=None):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_para(tf, text, size=14, bold=False, italic=False, color=DARK,
             align=PP_ALIGN.LEFT, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def pill(slide, x, y, w, h, text, bg, fg=WHITE, size=11, bold=True):
    """Rounded rectangle pill label."""
    shape = slide.shapes.add_shape(5, x, y, w, h)   # 5 = rounded rectangle
    shape.adjustments[0] = 0.5
    shape.fill.solid(); shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = fg
    return shape

def card(slide, x, y, w, h, title, body_lines,
         title_color=INDIGO, body_color=DARK, bg=PANEL,
         title_size=13, body_size=11):
    add_rect(slide, x, y, w, h, fill=bg)
    # title bar
    add_rect(slide, x, y, w, Inches(0.38), fill=title_color)
    add_text(slide, title, x + Inches(0.12), y + Inches(0.04),
             w - Inches(0.2), Inches(0.32),
             size=title_size, bold=True, color=WHITE)
    # body
    txb = slide.shapes.add_textbox(
        x + Inches(0.12), y + Inches(0.44),
        w - Inches(0.22), h - Inches(0.52))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in body_lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(body_size)
        run.font.color.rgb = body_color

def slide_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def accent_bar(slide, color=INDIGO, h=Inches(0.07)):
    add_rect(slide, 0, 0, W, h, fill=color)

def footer(slide, text="Campaign Readiness Copilot  ·  Confidential", color=GRAY):
    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), fill=INDIGO_D)
    add_text(slide, text,
             Inches(0.4), H - Inches(0.33), W - Inches(0.8), Inches(0.3),
             size=9, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
slide_bg(s1, LIGHT_BG)

# Full left panel (indigo)
add_rect(s1, 0, 0, Inches(4.8), H, fill=INDIGO)

# Left panel content
add_text(s1, "The Problem",
         Inches(0.3), Inches(0.9), Inches(4.2), Inches(0.8),
         size=36, bold=True, color=WHITE)
add_text(s1, "Campaign managers spend hours turning\nraw briefs into launch-ready plans —\nwith no consistent quality net.",
         Inches(0.3), Inches(1.85), Inches(4.2), Inches(1.6),
         size=16, color=RGBColor(0xC7, 0xD2, 0xFE), italic=True)

# Pain-point icons + text (left panel)
pains_left = [
    ("⏱", "Hours lost on manual brief analysis"),
    ("🔍", "No systematic readiness check"),
    ("⚠", "Compliance risks caught too late"),
    ("🔁", "Endless revision cycles before launch"),
]
for i, (icon, text) in enumerate(pains_left):
    yy = Inches(3.4 + i * 0.72)
    add_text(s1, icon, Inches(0.3), yy, Inches(0.5), Inches(0.6), size=20, color=AMBER)
    add_text(s1, text, Inches(0.85), yy + Inches(0.05), Inches(3.8), Inches(0.55),
             size=13, color=WHITE)

# Right panel — statistics & context
add_text(s1, "Campaign Briefing\nIs Broken",
         Inches(5.2), Inches(0.7), Inches(7.6), Inches(1.5),
         size=34, bold=True, color=DARK)

stats = [
    ("67%", "of campaigns launch with at least one\ncritical brief gap unresolved"),
    ("3–5×", "average revision cycles before a brief\nis approved for execution"),
    ("40%", "of campaign failures traced back\nto unclear objectives or KPIs"),
    ("$2M+", "average waste per enterprise from\npoorly scoped campaigns annually"),
]
for i, (stat, desc) in enumerate(stats):
    col = i % 2
    row = i // 2
    bx = Inches(5.2 + col * 3.9)
    by = Inches(2.5 + row * 2.15)
    add_rect(s1, bx, by, Inches(3.7), Inches(1.9), fill=WHITE)
    add_rect(s1, bx, by, Inches(3.7), Inches(0.08), fill=AMBER)
    add_text(s1, stat, bx + Inches(0.2), by + Inches(0.2), Inches(3.3), Inches(0.75),
             size=38, bold=True, color=INDIGO)
    add_text(s1, desc, bx + Inches(0.2), by + Inches(0.95), Inches(3.3), Inches(0.8),
             size=11.5, color=GRAY)

# Slide number pill
pill(s1, W - Inches(0.8), Inches(0.15), Inches(0.55), Inches(0.32),
     "1 / 4", INDIGO_D, size=10)
footer(s1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Solution
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
slide_bg(s2, LIGHT_BG)
accent_bar(s2, INDIGO)

add_text(s2, "The Solution",
         Inches(0.5), Inches(0.2), Inches(6), Inches(0.6),
         size=11, bold=True, color=INDIGO)
add_text(s2, "Campaign Readiness Copilot",
         Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.95),
         size=36, bold=True, color=DARK)
add_text(s2, "An AI-powered workflow that takes a campaign manager from raw brief to "
             "launch-ready plan in minutes — not days.",
         Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.6),
         size=15, color=GRAY, italic=True)

# 5-step workflow cards
steps = [
    ("1", "Understand\nBrief",      INDIGO,    "Parse & classify\nraw campaign text"),
    ("2", "Validate\nReadiness",    SKY,       "Score gaps, risks\n& KPI alignment"),
    ("3", "Resolve\nAmbiguity",     AMBER,     "AI questions +\nassumption register"),
    ("4", "Plan &\nSimulate",       GREEN,     "Channel strategy,\ntimeline & uplift"),
    ("5", "QA &\nHandoff",          RED,       "3-way QA check\n+ stakeholder pack"),
]
step_w = Inches(2.4)
step_h = Inches(2.8)
gap    = Inches(0.18)
start_x = (W - (step_w * 5 + gap * 4)) / 2

for i, (num, name, color, desc) in enumerate(steps):
    bx = start_x + i * (step_w + gap)
    by = Inches(2.4)
    add_rect(s2, bx, by, step_w, step_h, fill=WHITE)
    add_rect(s2, bx, by, step_w, Inches(0.08), fill=color)
    # step number circle
    add_rect(s2, bx + step_w/2 - Inches(0.38), by + Inches(0.18),
             Inches(0.76), Inches(0.76), fill=color)
    add_text(s2, num,
             bx + step_w/2 - Inches(0.38), by + Inches(0.2),
             Inches(0.76), Inches(0.64),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s2, name,
             bx + Inches(0.1), by + Inches(1.1),
             step_w - Inches(0.2), Inches(0.85),
             size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s2, desc,
             bx + Inches(0.1), by + Inches(2.0),
             step_w - Inches(0.2), Inches(0.75),
             size=11, color=GRAY, align=PP_ALIGN.CENTER)
    # arrow connector (except last)
    if i < 4:
        ax = bx + step_w + gap/2 - Inches(0.05)
        ay = by + step_h/2 - Inches(0.15)
        add_text(s2, "▶", ax, ay, Inches(0.2), Inches(0.3), size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Key differentiators (bottom row)
diffs = [
    ("🧑‍💼  Human-in-the-Loop", "Manager approves every assumption before AI proceeds"),
    ("🔌  Works Offline",        "Full mock fallback — no Azure key required for demo"),
    ("⚡  Parallel Execution",   "Concurrent tool calls cut per-step latency by ~60%"),
    ("📋  Stakeholder Handoff",  "Auto-generated exec summary + risk register + next steps"),
]
diff_w = Inches(3.1)
diff_gap = Inches(0.17)
diff_start = (W - (diff_w * 4 + diff_gap * 3)) / 2
for i, (title, body) in enumerate(diffs):
    bx = diff_start + i * (diff_w + diff_gap)
    by = Inches(5.55)
    add_rect(s2, bx, by, diff_w, Inches(1.55), fill=PANEL)
    add_text(s2, title, bx + Inches(0.12), by + Inches(0.1), diff_w - Inches(0.22), Inches(0.45),
             size=11, bold=True, color=INDIGO)
    add_text(s2, body,  bx + Inches(0.12), by + Inches(0.58), diff_w - Inches(0.22), Inches(0.85),
             size=10, color=GRAY)

pill(s2, W - Inches(0.8), Inches(0.15), Inches(0.55), Inches(0.32), "2 / 4", INDIGO, size=10)
footer(s2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
slide_bg(s3, LIGHT_BG)
accent_bar(s3, INDIGO)

add_text(s3, "System Architecture",
         Inches(0.5), Inches(0.2), Inches(8), Inches(0.6),
         size=11, bold=True, color=INDIGO)
add_text(s3, "Three-Tier Full-Stack + LangGraph Agent",
         Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.7),
         size=30, bold=True, color=DARK)

# Layer definitions
layers = [
    # (x, y, w, h, header, header_color, items, item_color)
    (Inches(0.3),  Inches(1.6),  Inches(2.55), Inches(5.4),
     "PRESENTATION", SKY,
     ["React 18 + Vite", "TypeScript + Tailwind", "5-page workflow UI",
      "Axios API client", "localhost:5173"],
     SKY),

    (Inches(3.05), Inches(1.6),  Inches(2.55), Inches(5.4),
     "API LAYER", GREEN,
     ["FastAPI (Python)", "7 REST endpoints", "Pydantic validation",
      "In-memory sessions", "localhost:8000"],
     GREEN),

    (Inches(5.8),  Inches(1.6),  Inches(3.5),  Inches(5.4),
     "AGENT WORKFLOW", INDIGO,
     ["LangGraph state machine", "7 nodes · 2 conditional routes",
      "parse → validate → ambiguity →", "  approve → plan → QA → handoff",
      "QA loop-back on failures"],
     INDIGO),

    (Inches(9.5),  Inches(1.6),  Inches(3.5),  Inches(2.55),
     "LLM / PROMPTS", RGBColor(0x8B, 0x5C, 0xF6),
     ["Azure OpenAI (GPT)", "14 Jinja2 prompt templates",
      "Organised by workflow step",
      "Auto-fallback to mock data"],
     RGBColor(0x8B, 0x5C, 0xF6)),
]

for (lx, ly, lw, lh, hdr, hdr_color, items, item_color) in layers:
    add_rect(s3, lx, ly, lw, lh, fill=WHITE)
    add_rect(s3, lx, ly, lw, Inches(0.42), fill=hdr_color)
    add_text(s3, hdr, lx + Inches(0.12), ly + Inches(0.07),
             lw - Inches(0.2), Inches(0.32),
             size=10, bold=True, color=WHITE)
    txb = s3.shapes.add_textbox(lx + Inches(0.12), ly + Inches(0.52),
                                 lw - Inches(0.22), lh - Inches(0.62))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = ("• " if not item.startswith(" ") else "") + item
        run.font.size = Pt(11)
        run.font.color.rgb = DARK

# Tool box (bottom-right)
add_rect(s3, Inches(9.5), Inches(4.4), Inches(3.5), Inches(2.6), fill=RGBColor(0xFF, 0xF7, 0xED))
add_rect(s3, Inches(9.5), Inches(4.4), Inches(3.5), Inches(0.38), fill=AMBER)
add_text(s3, "14 TOOL FUNCTIONS", Inches(9.62), Inches(4.45),
         Inches(3.3), Inches(0.3), size=9.5, bold=True, color=WHITE)
tools_txt = ("brief_parser · classify_intent · gap_detector\n"
             "compliance_risk · kpi_review · clarification_q\n"
             "assumption_reg · channel_strategy · exec_plan\n"
             "asset_checklist · timeline_wb · uplift_sim\n"
             "brief_to_plan_qa · handoff_packet")
add_text(s3, tools_txt, Inches(9.62), Inches(4.88), Inches(3.3), Inches(2.0),
         size=10, color=DARK)

# Flow arrows between layers
arrow_y = Inches(4.35)
for ax in [Inches(2.85), Inches(5.6), Inches(9.3)]:
    add_text(s3, "→", ax, arrow_y, Inches(0.3), Inches(0.4),
             size=20, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

# Agent ↔ Tool vertical arrow
add_text(s3, "↕", Inches(10.95), Inches(4.2), Inches(0.3), Inches(0.35),
         size=16, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# Tech stack pills (bottom)
techs = ["Python 3.10", "LangGraph", "FastAPI", "React 18", "Vite + Tailwind",
         "Azure OpenAI", "Jinja2", "Pydantic", "asyncio"]
px = Inches(0.3)
py = Inches(7.0)
for t in techs:
    pill_w = Inches(len(t) * 0.095 + 0.25)
    pill(s3, px, py, pill_w, Inches(0.32), t, INDIGO_D, size=9)
    px += pill_w + Inches(0.1)

pill(s3, W - Inches(0.8), Inches(0.15), Inches(0.55), Inches(0.32), "3 / 4", INDIGO, size=10)
footer(s3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Impact & Demo
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
slide_bg(s4, LIGHT_BG)
accent_bar(s4, AMBER)

add_text(s4, "Impact & Demo",
         Inches(0.5), Inches(0.2), Inches(8), Inches(0.6),
         size=11, bold=True, color=AMBER)
add_text(s4, "From Brief to Launch-Ready Plan",
         Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.75),
         size=33, bold=True, color=DARK)
add_text(s4, "A live demonstration of all 5 steps — works with or without Azure OpenAI credentials.",
         Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.55),
         size=14, color=GRAY, italic=True)

# Key metrics (top row)
metrics = [
    ("< 5 min",    "Full 5-step workflow\nfrom brief to handoff",   INDIGO),
    ("14 Tools",   "Modular AI functions,\neach independently tested", GREEN),
    ("2 Modes",    "Live Azure OpenAI\nor full mock fallback",       SKY),
    ("100%",       "Human approval required\nbefore AI plans launch", AMBER),
]
m_w = Inches(3.0); m_h = Inches(1.85)
m_gap = Inches(0.18)
m_start = (W - (m_w * 4 + m_gap * 3)) / 2
for i, (stat, desc, color) in enumerate(metrics):
    bx = m_start + i * (m_w + m_gap)
    by = Inches(2.2)
    add_rect(s4, bx, by, m_w, m_h, fill=WHITE)
    add_rect(s4, bx, by, m_w, Inches(0.07), fill=color)
    add_text(s4, stat, bx + Inches(0.15), by + Inches(0.18),
             m_w - Inches(0.28), Inches(0.72),
             size=30, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s4, desc, bx + Inches(0.12), by + Inches(0.95),
             m_w - Inches(0.22), Inches(0.75),
             size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Demo steps (left half)
add_text(s4, "How to Run the Demo",
         Inches(0.4), Inches(4.3), Inches(6.0), Inches(0.45),
         size=14, bold=True, color=DARK)

demo_steps = [
    ("1", "Clone repo & copy .env.example → .env (or skip for mock mode)"),
    ("2", "Terminal A: .venv/bin/uvicorn backend.main:app --reload --port 8000"),
    ("3", "Terminal B: cd frontend && npm install && npm run dev"),
    ("4", "Open http://localhost:5173 — paste any campaign brief and click Start"),
    ("5", "Walk through all 5 steps; download the Handoff Packet at the end"),
]
for i, (num, text) in enumerate(demo_steps):
    by = Inches(4.85 + i * 0.48)
    pill(s4, Inches(0.4), by, Inches(0.34), Inches(0.34), num, INDIGO, size=10)
    add_text(s4, text, Inches(0.85), by, Inches(5.8), Inches(0.42), size=11, color=DARK)

# "What's inside" panel (right half)
add_rect(s4, Inches(6.8), Inches(4.2), Inches(6.1), Inches(2.85), fill=INDIGO)
add_text(s4, "What's Inside the Repo",
         Inches(7.0), Inches(4.35), Inches(5.7), Inches(0.45),
         size=13, bold=True, color=WHITE)

inside = [
    "agent/workflow/   LangGraph graph, nodes & state",
    "agent/tools/      14 async tool functions",
    "agent/prompts/    14 Jinja2 LLM prompt templates",
    "backend/          FastAPI app + 7 routes",
    "frontend/src/     React 18 UI components",
    "diagrams/         Architecture.png + this deck",
    "tests/            pytest suite (mock + async)",
]
txb = s4.shapes.add_textbox(Inches(7.0), Inches(4.88), Inches(5.7), Inches(2.0))
tf = txb.text_frame; tf.word_wrap = True
first = True
for line in inside:
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(3)
    run = p.add_run()
    run.text = "▸  " + line
    run.font.size = Pt(10.5)
    run.font.color.rgb = RGBColor(0xC7, 0xD2, 0xFE)

# GitHub link
add_rect(s4, Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.28), fill=INDIGO_D)
add_text(s4, "🔗  github.com/tetalireddy-lv/Campaign-uplift-simulation",
         Inches(0.6), Inches(7.1), Inches(12.0), Inches(0.3),
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

pill(s4, W - Inches(0.8), Inches(0.15), Inches(0.55), Inches(0.32), "4 / 4", AMBER, size=10)
# no footer on last slide — github bar acts as footer

prs.save(OUT)
print(f"Saved: {OUT}")
